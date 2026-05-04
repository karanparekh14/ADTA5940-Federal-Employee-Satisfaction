"""
Build the Module 5 Final Presentation .pptx deck.
Run AFTER presentation_pipeline.py has produced figures_presentation/.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

HERE = Path(__file__).resolve().parent
FIGDIR = HERE / "figures_presentation"
OUT = HERE / "Module5_Final_Presentation_Parekh_Lee.pptx"

# ---------- Midnight Executive palette ----------
NAVY    = RGBColor(0x1E, 0x27, 0x61)
ICE     = RGBColor(0xCA, 0xDC, 0xFC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD    = RGBColor(0xFF, 0xC1, 0x07)   # accent for hook numbers
INK     = RGBColor(0x21, 0x29, 0x5C)   # darker navy for body text
MUTED   = RGBColor(0x6B, 0x73, 0x8D)   # secondary text

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

# ---------- Setup deck (16:9) ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=INK, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_left_accent(slide, color=NAVY, w=Inches(0.18)):
    return add_rect(slide, Inches(0), Inches(0), w, SH, color)


def add_page_number(slide, n, total):
    add_text(slide, SW - Inches(1.2), SH - Inches(0.45),
             Inches(1.0), Inches(0.3),
             f"{n} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer_label(slide, text):
    add_text(slide, Inches(0.4), SH - Inches(0.45),
             Inches(8), Inches(0.3),
             text, size=10, color=MUTED, align=PP_ALIGN.LEFT)


def title_block(slide, title, *, kicker=None, top=Inches(0.6)):
    if kicker:
        add_text(slide, Inches(0.6), top, Inches(11), Inches(0.4),
                 kicker.upper(), size=12, bold=True, color=NAVY,
                 font=FONT_HEAD)
        top = top + Inches(0.45)
    add_text(slide, Inches(0.6), top, Inches(12), Inches(1.0),
             title, size=32, bold=True, color=INK, font=FONT_HEAD,
             line_spacing=1.1)


# =====================================================================
# SLIDE 1 — TITLE (dark, full-bleed navy)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
# accent gold thin bar
add_rect(s, Inches(0.7), Inches(2.6), Inches(0.9), Inches(0.08), GOLD)
add_text(s, Inches(0.7), Inches(2.0), Inches(11), Inches(0.5),
         "ADTA 5940 — ANALYTICS CAPSTONE · SPRING 2026",
         size=13, bold=True, color=ICE, font=FONT_HEAD)
add_text(s, Inches(0.7), Inches(2.85), Inches(12), Inches(1.6),
         "Remote Work and\nFederal Employee Satisfaction",
         size=54, bold=True, color=WHITE, font=FONT_HEAD,
         line_spacing=1.05)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.6),
         "A view from 646,000 federal employees",
         size=22, color=ICE, font=FONT_HEAD)
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
         "Karan Parekh  ·  Chau Le",
         size=16, bold=True, color=WHITE, font=FONT_HEAD)
add_text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.3),
         "University of North Texas  ·  Dr. Denise Philpot",
         size=12, color=ICE, font=FONT_HEAD)


# =====================================================================
# SLIDE 2 — THE HOOK (one big number)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
add_text(s, Inches(0.6), Inches(0.5), Inches(8), Inches(0.5),
         "THE STARTING POINT", size=12, bold=True, color=NAVY)

# Giant 71%
add_text(s, Inches(0.5), Inches(1.6), Inches(7), Inches(3.5),
         "71%", size=240, bold=True, color=NAVY, font=FONT_HEAD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(5.3), Inches(7), Inches(0.6),
         "of federal employees say they're satisfied",
         size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.85), Inches(7), Inches(0.4),
         "Today's question: why those, and not the other 187,000?",
         size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Right side: distribution chart
s.shapes.add_picture(str(FIGDIR / "01_satisfaction_distribution.png"),
                     Inches(7.6), Inches(1.6), width=Inches(5.4))
add_footer_label(s, "Source: 2024 OPM Federal Employee Viewpoint Survey · n = 646,545")
add_page_number(s, 2, 13)


# =====================================================================
# SLIDE 3 — THE QUESTION (clean, almost-empty)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, ICE)
add_rect(s, Inches(0.6), Inches(3.2), Inches(0.18), Inches(1.1), NAVY)
add_text(s, Inches(0.95), Inches(2.7), Inches(11.8), Inches(0.5),
         "OUR RESEARCH QUESTION", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.95), Inches(3.2), Inches(11.6), Inches(2.6),
         "Does where you work — remote, hybrid, or on-site —\nactually change how satisfied you are?",
         size=36, bold=True, color=INK, font=FONT_HEAD, line_spacing=1.25)
add_page_number(s, 3, 13)


# =====================================================================
# SLIDE 4 — THE DATA (3 stat cards)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "The data", kicker="What we worked with")

card_y = Inches(2.4)
card_h = Inches(3.1)
card_w = Inches(3.7)
gap = Inches(0.4)
total_w = card_w * 3 + gap * 2
start_x = (SW - total_w) / 2

cards = [
    ("646K", "respondents",
     "Every federal employee who responded\nto the 2024 viewpoint survey."),
    ("36", "agencies",
     "From cabinet departments to small\nindependent agencies."),
    ("96", "questions",
     "Job satisfaction, engagement,\nwork-life balance, demographics."),
]
for i, (big, label, sub) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    add_rect(s, x, card_y, card_w, card_h, WHITE, line=ICE)
    # navy strip on top of each card
    add_rect(s, x, card_y, card_w, Inches(0.1), NAVY)
    add_text(s, x, card_y + Inches(0.4), card_w, Inches(1.4),
             big, size=72, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, card_y + Inches(1.7), card_w, Inches(0.4),
             label, size=18, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), card_y + Inches(2.2), card_w - Inches(0.6),
             Inches(1.0), sub, size=12, color=MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.3)

add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "Source: U.S. Office of Personnel Management — Public Release Data File",
         size=11, color=MUTED)
add_page_number(s, 4, 13)


# =====================================================================
# SLIDE 5 — APPROACH (4-step flow diagram)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "We built the model in four steps", kicker="Approach")

add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
         "Each step adds one ingredient and shows how much more it explains.",
         size=15, color=MUTED)

steps = [
    ("1", "Telework only",      "Where you work."),
    ("2", "+ Work-life balance","Do you have time for life?"),
    ("3", "+ Engagement",       "Are you energized by your work?"),
    ("4", "+ Demographics",     "Age, gender, tenure, supervisor."),
]
box_w = Inches(2.85)
box_h = Inches(3.0)
arrow_w = Inches(0.35)
total_w = box_w * 4 + arrow_w * 3
start_x = (SW - total_w) / 2
y = Inches(3.0)

for i, (num, head, sub) in enumerate(steps):
    x = start_x + i * (box_w + arrow_w)
    color = NAVY if i == len(steps) - 1 else INK
    add_rect(s, x, y, box_w, box_h, WHITE, line=ICE)
    add_rect(s, x, y, box_w, Inches(0.6), NAVY)
    add_text(s, x, y + Inches(0.07), box_w, Inches(0.5),
             num, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.25), y + Inches(0.95), box_w - Inches(0.5),
             Inches(1.2), head, size=18, bold=True, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.15)
    add_text(s, x + Inches(0.25), y + Inches(2.1), box_w - Inches(0.5),
             Inches(1.0), sub, size=13, color=MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.3)
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.05)
        ay = y + box_h / 2 - Inches(0.1)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                 Inches(0.25), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
        arr.line.fill.background()

add_page_number(s, 5, 13)


# =====================================================================
# SLIDE 6 — SURFACE FINDING (sets up surprise)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "At first glance, remote workers look more satisfied",
            kicker="Finding 1 — surface")

s.shapes.add_picture(str(FIGDIR / "02_telework_raw_gap.png"),
                     Inches(0.6), Inches(2.1), width=Inches(7.8))

# Right side callout
add_rect(s, Inches(8.8), Inches(2.4), Inches(4.0), Inches(3.7), ICE)
add_text(s, Inches(9.0), Inches(2.6), Inches(3.6), Inches(0.4),
         "THE GAP", size=11, bold=True, color=NAVY)
add_text(s, Inches(9.0), Inches(3.0), Inches(3.6), Inches(1.5),
         "0.40", size=72, bold=True, color=NAVY,
         font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.0), Inches(4.6), Inches(3.6), Inches(0.4),
         "points on the 5-point\nsatisfaction scale",
         size=13, color=INK, line_spacing=1.3)
add_text(s, Inches(9.0), Inches(5.4), Inches(3.6), Inches(0.6),
         "Routine remote (3.94)\nvs. required on-site (3.54)",
         size=11, color=MUTED, line_spacing=1.3)

add_footer_label(s, "FEVS 2024  ·  Q70 mean by Q91 telework status")
add_page_number(s, 6, 13)


# =====================================================================
# SLIDE 7 — THE TWIST: Simpson's Paradox
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "…but that's not the whole story",
            kicker="Finding 1 — the twist")

s.shapes.add_picture(str(FIGDIR / "04_simpsons_paradox.png"),
                     Inches(0.6), Inches(2.1), width=Inches(7.8))

# Right callout — flip explanation
add_rect(s, Inches(8.8), Inches(2.1), Inches(4.0), Inches(4.4), NAVY)
add_text(s, Inches(9.0), Inches(2.4), Inches(3.6), Inches(0.4),
         "WHAT HAPPENED", size=11, bold=True, color=GOLD)
add_text(s, Inches(9.0), Inches(2.9), Inches(3.6), Inches(1.4),
         "The telework\neffect flipped\nsign.",
         size=22, bold=True, color=WHITE, line_spacing=1.2)
add_text(s, Inches(9.0), Inches(4.6), Inches(3.6), Inches(2.0),
         "Once we account for engagement\nand work-life balance, being\nrequired on-site goes from\n−0.40 to slightly positive.\n\nThe original gap wasn't really\nabout where people work.",
         size=11, color=ICE, line_spacing=1.35)

add_footer_label(s, "Hierarchical OLS regression  ·  M1 vs. M4 telework coefficients")
add_page_number(s, 7, 13)


# =====================================================================
# SLIDE 8 — Headline R² progression
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "Engagement explains the variation, not telework",
            kicker="Finding 2 — what actually matters")

s.shapes.add_picture(str(FIGDIR / "03_r2_progression.png"),
                     Inches(0.5), Inches(2.0), width=Inches(8.4))

# Right side: 3 stat rows
rx = Inches(9.4)
rows = [
    ("2%",  "Telework alone"),
    ("41%", "+ Work-life balance"),
    ("60%", "+ Engagement"),
]
ry = Inches(2.4)
for big, label in rows:
    add_text(s, rx, ry, Inches(3.4), Inches(0.7),
             big, size=44, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(s, rx, ry + Inches(0.85), Inches(3.4), Inches(0.4),
             label, size=13, color=INK)
    ry += Inches(1.4)

add_footer_label(s, "R² of nested OLS models  ·  n = 519,284")
add_page_number(s, 8, 13)


# =====================================================================
# SLIDE 9 — Standardized coefficients
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "Engagement is ~7× larger than the next thing",
            kicker="Finding 2 — ranked drivers")

s.shapes.add_picture(str(FIGDIR / "05_standardized_coefs.png"),
                     Inches(0.5), Inches(2.0), width=Inches(7.4))

add_rect(s, Inches(8.3), Inches(2.0), Inches(4.6), Inches(4.7), ICE)
add_text(s, Inches(8.5), Inches(2.2), Inches(4.2), Inches(0.4),
         "READING THIS CHART", size=11, bold=True, color=NAVY)
add_text(s, Inches(8.5), Inches(2.7), Inches(4.2), Inches(4.0),
         "Each bar is the standardized effect on satisfaction.\n\nEngagement: β = 0.71\nNext closest: β = 0.10\n\nWhere you work, gender, age, tenure — barely register.",
         size=12, color=INK, line_spacing=1.4)

add_footer_label(s, "Standardized OLS · Model 4")
add_page_number(s, 9, 13)


# =====================================================================
# SLIDE 10 — Bonus: Attrition  (uses Chau's notebook RF feature importance)
# =====================================================================
NB_FIGS = HERE / "figures_from_notebook"

s = prs.slides.add_slide(BLANK)
add_left_accent(s, color=GOLD)  # gold accent to mark "bonus"
add_text(s, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "BONUS — SECONDARY ANALYSIS", size=12, bold=True, color=GOLD)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(1.0),
         "Side question: who's likely to leave?",
         size=32, bold=True, color=INK, font=FONT_HEAD)

# Left: three stat callouts
lx = Inches(0.6)
ly = Inches(2.3)
add_text(s, lx, ly, Inches(4.0), Inches(1.3),
         "33%", size=80, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, lx, ly + Inches(1.4), Inches(4.0), Inches(0.5),
         "of federal employees are considering leaving",
         size=14, color=INK, line_spacing=1.3)

add_text(s, lx, ly + Inches(2.3), Inches(4.0), Inches(0.7),
         "AUC 0.77", size=36, bold=True, color=NAVY)
add_text(s, lx, ly + Inches(3.0), Inches(4.0), Inches(0.5),
         "XGBoost classifier (recall 64% on leavers)",
         size=12, color=MUTED, line_spacing=1.3)

# Right: feature importance chart from notebook (cell 34 — now XGBoost)
s.shapes.add_picture(
    str(NB_FIGS / "cell34_xgboost_feature_importance.png"),
    Inches(5.2), Inches(2.1), width=Inches(7.4))

# Translation key at bottom — decode the variable codes
add_text(s, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.3),
         "Top 3 drivers: Job satisfaction (Q70), Leadership (Q57–Q65 mean), Pay satisfaction (Q71).",
         size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.25),
         "Demographics (DFEDTEN_C = tenure 20+, DSEX_B = female, etc.) barely register.",
         size=9, color=MUTED, align=PP_ALIGN.CENTER)
add_page_number(s, 10, 13)


# =====================================================================
# SLIDE 11 — Bonus continued: pay → attrition (cell 16, corrected)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s, color=GOLD)
add_text(s, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "BONUS — THE PAY GRADIENT", size=12, bold=True, color=GOLD)
add_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(1.0),
         "And the relationship is monotonic",
         size=32, bold=True, color=INK, font=FONT_HEAD)

# Chart on the left — tighter so the source line at bottom doesn't collide
s.shapes.add_picture(
    str(NB_FIGS / "cell16_attrition_by_pay_CORRECTED.png"),
    Inches(0.6), Inches(2.0), width=Inches(7.6))

# Right side: gradient narrative
add_rect(s, Inches(8.6), Inches(2.0), Inches(4.3), Inches(4.7), ICE)
add_text(s, Inches(8.8), Inches(2.2), Inches(3.9), Inches(0.4),
         "FROM RED TO GREEN", size=11, bold=True, color=NAVY)
add_text(s, Inches(8.8), Inches(2.7), Inches(3.9), Inches(1.0),
         "65%", size=44, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, Inches(8.8), Inches(3.6), Inches(3.9), Inches(0.7),
         "attrition rate among employees very dissatisfied with pay",
         size=11, color=INK, line_spacing=1.3)

add_text(s, Inches(8.8), Inches(4.6), Inches(3.9), Inches(1.0),
         "20%", size=44, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, Inches(8.8), Inches(5.5), Inches(3.9), Inches(0.7),
         "attrition rate among those very satisfied with pay",
         size=11, color=INK, line_spacing=1.3)

add_footer_label(s, "Source: Attrition_Pipeline_Cleaned.ipynb · cell 16")
add_page_number(s, 11, 13)


# =====================================================================
# SLIDE 12 — Recommendations (3 takeaway rows)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_left_accent(s)
title_block(s, "Three takeaways for federal agencies", kicker="So what")

takeaways = [
    ("01", "Don't fight over telework policy.",
     "It's a small lever. Engagement is the big one."),
    ("02", "Invest in engagement.",
     "Manager quality, mission clarity, recognition — that's the 60%."),
    ("03", "Watch attrition signals.",
     "Low pay satisfaction × low engagement = high flight risk."),
]
y = Inches(2.4)
for num, head, body in takeaways:
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(0.7), y, Inches(0.95), Inches(0.95))
    circle.fill.solid(); circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    add_text(s, Inches(0.7), y, Inches(0.95), Inches(0.95),
             num, size=20, bold=True, color=WHITE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.0), y + Inches(0.05), Inches(11), Inches(0.5),
             head, size=22, bold=True, color=INK, font=FONT_HEAD)
    add_text(s, Inches(2.0), y + Inches(0.6), Inches(11), Inches(0.6),
             body, size=14, color=MUTED, line_spacing=1.3)
    y += Inches(1.45)

add_page_number(s, 12, 13)


# =====================================================================
# SLIDE 13 — Limitations + Thank You (dark close)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, Inches(0.7), Inches(0.85), Inches(0.9), Inches(0.08), GOLD)
add_text(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.4),
         "LIMITATIONS & WHAT WE'D DO NEXT", size=12, bold=True,
         color=ICE, font=FONT_HEAD)

# Three brief limitations
lims = [
    ("Cross-sectional", "We see association, not causation."),
    ("Self-reported",   "Same person rates engagement and satisfaction."),
    ("Single year",     "2024 only. Replicate across 2022–24 next."),
]
y = Inches(1.5)
for h, sub in lims:
    add_text(s, Inches(0.7), y, Inches(5), Inches(0.4),
             h, size=18, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(0.7), y + Inches(0.45), Inches(7), Inches(0.4),
             sub, size=13, color=ICE, line_spacing=1.3)
    y += Inches(1.0)

# Closing line — prominent
add_rect(s, Inches(0.7), Inches(5.0), Inches(0.18), Inches(1.5), GOLD)
add_text(s, Inches(1.0), Inches(5.0), Inches(11.5), Inches(0.7),
         "Telework gets the headlines.",
         size=28, color=ICE, font=FONT_HEAD, line_spacing=1.1)
add_text(s, Inches(1.0), Inches(5.7), Inches(11.5), Inches(0.7),
         "Engagement gets the results.",
         size=28, bold=True, color=WHITE, font=FONT_HEAD, line_spacing=1.1)

add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Thank you  ·  Questions?",
         size=14, color=ICE, font=FONT_HEAD)
add_text(s, SW - Inches(2.5), SH - Inches(0.45), Inches(2.0), Inches(0.3),
         "Karan Parekh · Chau Le", size=10, color=ICE,
         align=PP_ALIGN.RIGHT)


# =====================================================================
# BACKUP SLIDES (Q&A only — not in main flow)
# =====================================================================
def backup_slide(prs, title, kicker, image_path, narrative, idx, total=5,
                 chart_width=7.0, chart_y=2.1):
    """Backup slide. chart_width must keep the picture inside 7.5" slide height
       given the picture's intrinsic aspect ratio."""
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, Inches(0.18), SH, MUTED)
    add_text(s, Inches(0.6), Inches(0.4), Inches(8), Inches(0.4),
             f"BACKUP {idx} / {total} — {kicker}",
             size=11, bold=True, color=MUTED)
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             title, size=26, bold=True, color=INK, font=FONT_HEAD,
             line_spacing=1.1)
    s.shapes.add_picture(str(image_path),
                         Inches(0.6), Inches(chart_y),
                         width=Inches(chart_width))
    # Right side narrative panel — sized to fill remaining horizontal space
    panel_x = Inches(0.6 + chart_width + 0.4)
    panel_w = SW - panel_x - Inches(0.4)
    add_rect(s, panel_x, Inches(chart_y), panel_w, Inches(4.5), ICE)
    add_text(s, panel_x + Inches(0.2), Inches(chart_y + 0.2),
             panel_w - Inches(0.4), Inches(4.2),
             narrative, size=12, color=INK, line_spacing=1.4)
    add_footer_label(s, "Q&A backup — not part of main 15-min flow")
    return s


backup_slide(prs,
    "Pay satisfaction vs. leaving (boxplot)",
    "Q71 BOXPLOT",
    NB_FIGS / "cell14_import_matplotlib_pyplot_as_plt.png",
    "Stayers (0): median pay\nsatisfaction = 4.\n\nLeavers (1): median pay\nsatisfaction = 3.\n\nThe full distribution\nshifts down a full point\nfor those considering\nleaving.",
    1, chart_width=5.5)

backup_slide(prs,
    "Attrition varies a lot by agency",
    "TOP 20 AGENCIES",
    NB_FIGS / "cell25_colors_a0c4ff_len_attrition_by_dept.png",
    "Highest attrition: AF, DD,\nAR (~40-45%).\n\nLowest: GS, EP, CM\n(~22-25%).\n\nTwo-fold variation\nsuggests agency-level\nculture matters — not\njust individual factors.",
    2, chart_width=6.0)

backup_slide(prs,
    "Less satisfied agencies have higher attrition",
    "AGENCY SCATTER",
    NB_FIGS / "cell27_sns_scatterplot.png",
    "Each dot is one agency.\n\nClear inverse pattern:\nhigher mean job\nsatisfaction ↔ lower\nattrition rate.\n\nThe outlier in the\nupper-left is a small,\nlow-satisfaction agency\nwith high churn.",
    3, chart_width=5.5)

backup_slide(prs,
    "Department satisfaction heatmap",
    "Q70 / Q71 / LEADERSHIP BY AGENCY",
    NB_FIGS / "cell22_dept_scores_df_groupby_agency.png",
    "Mean scores across\n36 agencies for:\n• Job satisfaction (Q70)\n• Pay satisfaction (Q71)\n• Leadership composite\n\nDarker = higher score.\nMost agencies cluster\naround 3.7-4.0 on the\n5-point scale.",
    4, chart_width=5.5)

backup_slide(prs,
    "Model performance — XGBoost confusion matrix",
    "LR vs XGBOOST",
    NB_FIGS / "cell33_xgboost_confusion_matrix.png",
    "XGBoost\nconfusion matrix\non 20% holdout\n(n = 126,460).\n\nAUC: 0.771\nRecall on leavers: 64%\n(vs. 46% for RF —\nscale_pos_weight\nbalancing helped).\n\nLogistic Regression\nbaseline AUC: 0.764.",
    5, chart_width=5.5)


# ---------- Save ----------
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
